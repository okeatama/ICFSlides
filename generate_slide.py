from datetime import datetime, timedelta
from pptx import Presentation
import requests
from bs4 import BeautifulSoup
import os
from urllib.parse import urljoin
import json
from constants import *
import csv
from time import sleep
import re

# Target website


with open("JSON/calendar.json","r") as file:
    calendar = json.load(file)

with open("JSON/updated_mass_readings.json", "r") as file:
    mass_readings = json.load(file)

# gets date of next sunday
def get_next_sunday():
    today = datetime.now()
    days_ahead = (6 - today.weekday()) % 7
    if days_ahead == 0:
        days_ahead = 7  # skip to next week if today is Sunday
    return today + timedelta(days=days_ahead)

def create_warning_slide(prs, id_verses, en_verses):
    warning = prs.slides.add_slide(prs.slide_layouts[WARNING_LAYOUT])
    warning.placeholders[WARNING_TEXT].text = f"ID: {id_verses}\n\nEN:{en_verses}"

def create_birthday_slide(prs, mass_date):
    with open("resources/data.csv", "r", encoding='utf-8-sig') as file:
        csv_reader = csv.DictReader(file)
        durs = [timedelta(days=i) for i in range(7)]
        # dur = timedelta(days=6)
        from_date = mass_date - timedelta(days=6)

        # acceptable values
        # date_range = [str(i) for i in range(from_date.day, mass_date.day + 1)]
        date_range = {str((mass_date - dur).day) for dur in durs}
        month_range = set()
        month_range.add(str(from_date.month))
        month_range.add(str(mass_date.month))

        acceptable_active = ["Active", "Not Sure"]

        lines = []

        for row in csv_reader:
            # if row["Month"] in month_range and row["Date"] in date_range and row["Active"] in acceptable_active:
            if row["Month"] in month_range and row["Date"] in date_range:
                # birthday in range
                # full_name = f"{row['First Name']} {row['Middle Name']} {row['Last Name']}"
                full_name = f"{row['Full Name'].strip()}"
                line = f"{row['Date']} {NUM_TO_MONTH_ID[int(row['Month'])]}: {full_name}"
                lines.append(line)
    
    bday = prs.slides.add_slide(prs.slide_layouts[BIRTHDAY_LAYOUT])
    final_text = f"Bagi Mereka yang berulang tahun pada, {from_date.day} - {mass_date.day} {NUM_TO_MONTH_ID[mass_date.month]}:"
    final_text = final_text + '\n\n' + '\n'.join(lines)
    bday.placeholders[BIRTHDAY_PLACEHOLDER].text = final_text

def create_first_reading_slide(prs, content_id, content_en, bible_loc):
    reading = prs.slides.add_slide(prs.slide_layouts[FIRST_READING_LAYOUT])
    if len(content_id) > SLIDE_CONTENT_THRESHOLD:
        create_first_reading_slide(prs, content_id[SLIDE_CONTENT_THRESHOLD + 1:], "", bible_loc)
        content_id = content_id[:SLIDE_CONTENT_THRESHOLD + 1]
    reading.placeholders[FIRST_READING_INDO].text = content_id
    reading.placeholders[FIRST_READING_EN].text = content_en
    reading.placeholders[FIRST_READING_VERSES].text = bible_loc
    

def create_second_reading_slide(prs, content_id, content_en, bible_loc):
    reading = prs.slides.add_slide(prs.slide_layouts[SECOND_READING_LAYOUT])
    if len(content_id) > SLIDE_CONTENT_THRESHOLD:
        create_second_reading_slide(prs, content_id[SLIDE_CONTENT_THRESHOLD + 1:], "", bible_loc)
        content_id = content_id[:SLIDE_CONTENT_THRESHOLD + 1]
    reading.placeholders[SECOND_READING_INDO].text = content_id
    reading.placeholders[SECOND_READING_EN].text = content_en
    reading.placeholders[SECOND_READING_VERSES].text = bible_loc


def create_gospel_slide(prs, content_id, content_en, bible_loc):
    gospel = prs.slides.add_slide(prs.slide_layouts[GOSPEL_LAYOUT])
    if len(content_id) > SLIDE_CONTENT_THRESHOLD: 
        create_gospel_slide(prs, content_id[SLIDE_CONTENT_THRESHOLD + 1:], "", bible_loc)
        content_id = content_id[:SLIDE_CONTENT_THRESHOLD + 1]
    gospel.placeholders[GOSPEL_INDO].text = content_id
    gospel.placeholders[GOSPEL_EN].text = content_en
    gospel.placeholders[GOSPEL_VERSES].text = bible_loc

def create_resp_psalm_slide(prs, bible_loc, text_id, text_en):
    responsorial_psalm = prs.slides.add_slide(prs.slide_layouts[PSALM_LAYOUT])
    responsorial_psalm.placeholders[PSALM_VERSES].text = bible_loc
    responsorial_psalm.placeholders[PSALM_IMAGE].insert_picture(f"resources/{FILENAMES[0]}")
    responsorial_psalm.placeholders[PSALM_TEXT_EN].text = text_en
    responsorial_psalm.placeholders[PSALM_TEXT_INDO].text = text_id

def create_gospel_acclamation_slide(prs, verse_id, verse_en):
    gospel_acclamation = prs.slides.add_slide(prs.slide_layouts[GOSPEL_ACCLAMATION_LAYOUT])
    gospel_acclamation.placeholders[GOSPEL_ACCLAMATION_IMAGE].insert_picture(f"resources/{FILENAMES[1]}")
    gospel_acclamation.placeholders[GOSPEL_ACCLAMATION_TEXT_INDO].text = verse_id
    gospel_acclamation.placeholders[GOSPEL_ACCLAMATION_TEXT_EN].text = verse_en

def download_images_and_get_verse(mass_date):
    # returns ([mazmurs], ayat): ([String], String)
    # download responsorial psalm and gospel acclamation images
    
    mass_title_id = calendar[mass_date.month][mass_date.strftime("%Y-%m-%d")]

    for y in tries:
        url = f"{LAGUMISA_BASE_URL}{y}-{mass_title_id}"

        # try this url, and search for images in edisibaru
        response = requests.get(url)
        soup = BeautifulSoup(response.text, "html.parser")

        edisibaru = soup.find("div", id="edisibaru")

        if not edisibaru:
            continue

        # 1: since there is always an image before the 2 we want
        imgs = edisibaru.find_all("img")[1:]

        for i, img_tag in enumerate(imgs):
            img_url = urljoin(url, img_tag.get("src"))
            img_name = os.path.basename(img_url.split("?")[0])  # remove query params
            img_data = requests.get(img_url).content
            with open(f"resources/{FILENAMES[i]}", "wb") as f:
                f.write(img_data)
            print(f"Downloaded {FILENAMES[i]}")

        # get mazmur and ayat
        # [2:] since there is `Mazmur (oleh pemazmur):\n\n` and `Ayat (oleh solis):\n\n`
        presyair = edisibaru.find_all("pre", class_="presyair")
        mazmur_raw = "\n".join(presyair[0].text.split('\n')[2:])
        mazmur_split = mazmur_raw.split("\n\n")
        # [3:] to remove first 3 letters, which are 1. , 2. , 3. , etc
        # there's 3 spaces after each \n, which is annoying so remove it
        mazmur = [m[3:].replace("\n   ", " ") for m in mazmur_split] 

        ayat = "\n".join(presyair[1].text.split('\n')[2:])

        mazmur_verse = edisibaru.find("a", class_="aayat").text
        return (mazmur, mazmur_verse, ayat)

def extract_universalis(mass_date):
    # scrape from https://universalis.com/mass.htm
    url = f'{UNIVERSALIS_BASE_URL}{mass_date.strftime("%Y%m%d")}/mass.htm'

    response = requests.get(url)
    soup = BeautifulSoup(response.text, "html.parser")

    texts = soup.find("div", id="texts")
    content = texts.find("div").get_text()
    table = texts.find_all("th", align="right") # theres 5 usually, first verse, psalm, second, gospel acclamation, gospel

    # verses
    first_verse = table[0].get_text()
    psalm_loc = table[1].get_text()
    second_verse = table[2].get_text()
    gospel_verse = table[4].get_text()

    first_reading = ""
    second_reading = ""
    gospel = ""
    # signifies if they have collected first, second and gospel
    flags = [False, False, False] # first reading second reading, gospel

    line_iter = iter(content.split('\n'))
    sentinel = object()
    while True:
        line = next(line_iter, sentinel)
        if line is sentinel:
            # out of iteration
            break

        if "First reading" in line and not flags[0]:
            readings = [] # array of strings of line of verses
            while True:
                next_line = next(line_iter, sentinel)
                if (next_line is sentinel) or next_line.strip().startswith("How to listen"):
                    # break out of the iteration
                    break
                readings.append(next_line.strip())
            
            first_reading = "\n".join(readings)
            flags[0] = True
        elif "Second reading" in line and not flags[1]:
            readings = []
            while True:
                next_line = next(line_iter, sentinel)
                if (next_line is sentinel) or next_line.strip().startswith("Gospel Acclamation"):
                    # break out of the iteration
                    break
                readings.append(next_line.strip())
            second_reading = "\n".join(readings)
            flags[1] = True
        elif "Gospel Acclamation" in line:
            # skip over Gospel Acclamation, since this word contains "Gospel"
            continue
        elif "Gospel" in line and not flags[2]:
            readings = []
            while True:
                next_line = next(line_iter, sentinel)
                if (next_line is sentinel) or next_line.strip().startswith("The responsorial psalms"):
                    # break out of the iteration
                    break
                readings.append(next_line.strip())
            gospel = "\n".join(readings)
            flags[2] = True
        
        if all(flags):
            # all three has been extracted, no need to go any further
            break
    
    return (first_reading, first_verse, psalm_loc, second_reading, second_verse, gospel, gospel_verse)

def extract_imankatolik(mass_date):
    # scrape from https://www.imankatolik.or.id/kalender.php
    url = f"{IMANKATOLIK_BASE_URL}/kalender.php?b={mass_date.month}&t={mass_date.year}"

    response = requests.get(url)
    soup = BeautifulSoup(response.text, "html.parser")

    mass_date_section = soup.find("a", href=f"/kalender/{mass_date.day}{NUM_TO_MONTH_ID[mass_date.month]}.html")

    tr = mass_date_section.find_parent("tr")
    readings = tr.next_sibling.next # readings should start in sibling
    urls = []

    for i in readings.children:
        # get the urls leading to the readings
        if i.text == "." or i.text == "<br/>":
            # additional stuff we don't need (hopefully fingers crossed)
            break
        elif i.name != 'a':
            # find for <a> tags, anything else ignore it
            continue

        urls.append(i.attrs["href"])

    texts = []
    for url in urls:
        # go to each urls and scrape contents
        # give some time delay so server isnt overloaded
        sleep(1)
        new_url = f"{IMANKATOLIK_BASE_URL}{url}"
        response = requests.get(new_url)
        soup = BeautifulSoup(response.text, "html.parser")

        temp = []
        for text in soup.find_all("td", width="95%"):
            temp.append(text.text)
        texts.append("\n".join(temp))
    
    reading_verses = readings.text.split("\r\n", 1)[0]
    # return the texts for each urls, order should be [first reading, psalm (not needed), second reading, gospel]
    # and second element is all the verses in indo
    return texts, reading_verses

def get_en_psalms(verses):
    # verses is expected to be like MAZMUR 143:1-2.5-6.7ab.8ab.10;R:1a

    # removes MAZMUR, can assume its from PSALM, and remove last bits after ; which we don't need
    verses = verses.split(" ", 1)[1].split(";")[0] 
    # 143:1-2.5-6.7ab.8ab.10
    # remove any letters
    verses = verses.replace("a", "").replace("b", "").replace("c", "").replace("d", "")

    # 143:1-2.5-6.7.8.10
    verses_split = verses.split(":")
    verse_list = []
    chapter = None
    if len(verses_split) == 2:
        # only 1 chapter hurray simple
        chapter = verses_split[0]

        # verse_split[1] == "1-2.5-6.7.8.10"
        verse_list = expand_ranges(verses_split[1])
    else:
        print("EXPLODE CUZ OF DIFFERENT CHAPTERS IN PSALMS KABOOM")
    
    res_verses = []
    with open("JSON/psalm.json","r") as file:
        psalms = json.load(file)

        for verse in verse_list:
            res_verses.append(psalms[chapter][str(verse)])

    return "\n".join(res_verses)

def expand_ranges(s):
    # this is vibe coded
    # Split by dot or comma
    parts = re.split(r'[.,;]', s)
    numbers = []

    for part in parts:
        if '-' in part:
            start, end = map(int, part.split('-'))
            numbers.extend(range(start, end + 1))
        else:
            numbers.append(int(part))
    return numbers

def main():

    mass_date = get_next_sunday()

    # scrape from lagumisa
    mazmur, mazmur_verse, ayat = download_images_and_get_verse(mass_date)

    # scrape from sources
    first_reading, first_verse, psalm_loc, second_reading, second_verse, gospel, gospel_verse = extract_universalis(mass_date)
    texts_id, id_verses = extract_imankatolik(mass_date)
    psalms = get_en_psalms(mazmur_verse)

    # create the presentation
    prs = Presentation("resources/template.pptx")
    create_warning_slide(prs, id_verses, f"{first_verse};{second_verse};{gospel_verse}.")
    create_birthday_slide(prs, mass_date)
    create_first_reading_slide(prs, texts_id[0], first_reading, first_verse)

    for text in mazmur:
        create_resp_psalm_slide(prs, psalm_loc, text, psalms)

    create_second_reading_slide(prs, texts_id[2], second_reading, second_verse)

    create_gospel_acclamation_slide(prs, ayat, "Translate it yourself please I have no idea where to find the source :)")

    create_gospel_slide(prs, texts_id[3], gospel, gospel_verse)
    pptxfilename = mass_date.strftime("%Y%m%d")
    prs.save(f"{pptxfilename}.pptx")

if __name__ == "__main__":
    main()



""" 
old code using the bible here

def extract_mass_reading(mass_reading):
    first_readings = []
    first_verses = []
    second_readings = []
    second_verses = []
    gospel = []
    gospel_verses = []
    psalm_loc = ""

    for section in mass_reading:
        if section["header"] == "Reading 1":
            for r in section["readings"]:
                first_readings.append(r["text"])
                first_verses.append(r["verses"][0]["text"])
        
        elif section["header"] == "Reading 2":
            for r in section["readings"]:
                second_readings.append(r["text"])
                second_verses.append(r["verses"][0]["text"])
        elif section["header"] == "Responsorial Psalm":
            psalm_loc = section["readings"][0]["verses"][0]["text"]
        elif section["header"] == "Gospel":
            for r in section["readings"]:
                gospel.append(r["text"])
                gospel_verses.append(r["verses"][0]["text"])
    
    return (first_readings, first_verses, psalm_loc, second_readings, second_verses, gospel, gospel_verses)

mass_reading = None

# find mass_reading for mass date
for i, r in enumerate(mass_readings):
    if r["date"] == mass_date.strftime("%Y-%m-%d"):
        mass_reading = mass_readings[i]["sections"]
        break

if not mass_reading:
    raise Exception(f"Couldn't find date {mass_date.strftime('%Y-%m-%d')}")
# now do some magic to extract readings, psalm and gospel
first_readings, first_verses, psalm_loc, second_readings, second_verses, gospel, gospel_verses = extract_mass_reading(mass_reading)
"""

